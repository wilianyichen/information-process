from __future__ import annotations

from pathlib import Path
from typing import Any

from infoproc.adapters.base import InputAdapter
from infoproc.models import InputKind


DOCUMENT_EXTENSIONS = {".pdf", ".doc", ".docx", ".ppt", ".pptx"}
PRESENTATION_EXTENSIONS = {".ppt", ".pptx"}


class DocumentInputAdapter(InputAdapter):
    kind = InputKind.DOCUMENT

    def supports(self, path: Path) -> bool:
        return path.suffix.lower() in DOCUMENT_EXTENSIONS

    def probe(self, path: Path) -> dict[str, Any]:
        suffix = path.suffix.lower()
        detected_kind = InputKind.PRESENTATION if suffix in PRESENTATION_EXTENSIONS else InputKind.DOCUMENT
        payload: dict[str, Any] = {
            "input_path": str(path),
            "detected_kind": detected_kind.value,
            "extension": suffix,
            "requires_conversion": suffix in {".doc", ".ppt"},
            "normalized_extension": ".docx" if suffix == ".doc" else ".pptx" if suffix == ".ppt" else suffix,
        }
        if suffix == ".pdf":
            payload.update(_probe_pdf(path))
        elif suffix == ".docx":
            payload.update(_probe_docx(path))
        elif suffix == ".pptx":
            payload.update(_probe_pptx(path))
        return payload

    def normalize(self, path: Path, output_audio_path: Path) -> dict[str, str]:
        return {
            "skipped": True,
            "reason": "document normalization is handled by DocumentExtractionService",
        }


def _probe_pdf(path: Path) -> dict[str, Any]:
    try:
        from pypdf import PdfReader
    except Exception:
        return {}
    try:
        reader = PdfReader(str(path))
        return {"page_count": len(reader.pages)}
    except Exception:
        return {}


def _probe_docx(path: Path) -> dict[str, Any]:
    try:
        from docx import Document
    except Exception:
        return {}
    try:
        document = Document(str(path))
        text_chunks: list[str] = []
        paragraph_count = 0
        for paragraph in document.paragraphs:
            paragraph_count += 1
            if paragraph.text.strip():
                text_chunks.append(paragraph.text.strip())
        return {
            "paragraph_count": paragraph_count,
            "characters": sum(len(item) for item in text_chunks),
        }
    except Exception:
        return {}


def _probe_pptx(path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except Exception:
        return {}
    try:
        presentation = Presentation(str(path))
        characters = 0
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "") or ""
                characters += len(text.strip())
        return {
            "slide_count": len(presentation.slides),
            "characters": characters,
        }
    except Exception:
        return {}
