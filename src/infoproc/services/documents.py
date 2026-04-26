from __future__ import annotations

from pathlib import Path
import shutil
import tempfile

from infoproc.config import AppConfig
from infoproc.utils import run_command


class DocumentExtractionService:
    def __init__(self, config: AppConfig) -> None:
        self.config = config

    def convert_legacy_office(self, source_path: Path, output_path: Path) -> dict[str, str]:
        converter = self._resolve_office_converter()
        if converter is None:
            raise RuntimeError(
                "Legacy .doc/.ppt conversion requires LibreOffice headless. "
                f"Configured converter not found: {self.config.document.office_converter}"
            )
        output_path.parent.mkdir(parents=True, exist_ok=True)
        result = run_command(
            [
                converter,
                "--headless",
                "--convert-to",
                output_path.suffix.lstrip("."),
                "--outdir",
                str(output_path.parent),
                str(source_path),
            ],
            timeout=1800,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice conversion failed for {source_path.name}: {result.stderr.strip() or result.stdout.strip()}"
            )
        converted = output_path.parent / f"{source_path.stem}{output_path.suffix}"
        if not converted.exists():
            raise RuntimeError(f"LibreOffice did not produce expected output: {converted}")
        if converted != output_path:
            output_path.unlink(missing_ok=True)
            converted.replace(output_path)
        return {
            "converter": converter,
            "normalized_path": str(output_path),
        }

    def extract_text(self, source_path: Path) -> tuple[dict[str, str | int], str]:
        suffix = source_path.suffix.lower()
        if suffix == ".pdf":
            return self._extract_pdf_text(source_path)
        if suffix == ".docx":
            return self._extract_docx_text(source_path)
        if suffix == ".pptx":
            return self._extract_pptx_text(source_path)
        raise RuntimeError(f"Unsupported document extraction format: {suffix}")

    def _extract_pdf_text(self, source_path: Path) -> tuple[dict[str, str | int], str]:
        engine = self.config.document.pdf_engine.lower()
        if engine == "pdftotext":
            text = self._extract_pdf_with_pdftotext(source_path)
            if text.strip():
                return {"engine": "pdftotext"}, _normalize_document_text(text)
            text = self._extract_pdf_with_pypdf(source_path)
            return {"engine": "pypdf_fallback"}, _normalize_document_text(text)

        text = self._extract_pdf_with_pypdf(source_path)
        if text.strip():
            return {"engine": "pypdf"}, _normalize_document_text(text)
        fallback = self._extract_pdf_with_pdftotext(source_path)
        if fallback.strip():
            return {"engine": "pdftotext_fallback"}, _normalize_document_text(fallback)
        return {"engine": "pypdf"}, ""

    def _extract_pdf_with_pypdf(self, source_path: Path) -> str:
        try:
            from pypdf import PdfReader
        except Exception as exc:
            raise RuntimeError("pypdf is required for PDF text extraction.") from exc
        reader = PdfReader(str(source_path))
        parts: list[str] = []
        for page in reader.pages:
            parts.append((page.extract_text() or "").strip())
        return "\n\n".join(part for part in parts if part)

    def _extract_pdf_with_pdftotext(self, source_path: Path) -> str:
        binary = shutil.which("pdftotext")
        if binary is None:
            return ""
        with tempfile.TemporaryDirectory(prefix="infoproc-pdf-") as tmp:
            output_path = Path(tmp) / "out.txt"
            result = run_command([binary, "-layout", str(source_path), str(output_path)], timeout=1800)
            if result.returncode != 0 or not output_path.exists():
                return ""
            return output_path.read_text(encoding="utf-8", errors="ignore")

    def _extract_docx_text(self, source_path: Path) -> tuple[dict[str, str | int], str]:
        try:
            from docx import Document
        except Exception as exc:
            raise RuntimeError("python-docx is required for DOCX text extraction.") from exc
        document = Document(str(source_path))
        lines: list[str] = []
        paragraph_count = 0
        for paragraph in document.paragraphs:
            paragraph_count += 1
            text = paragraph.text.strip()
            if text:
                lines.append(text)
        for table in document.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    lines.append(row_text)
        return {
            "engine": "python-docx",
            "paragraph_count": paragraph_count,
        }, _normalize_document_text("\n".join(lines))

    def _extract_pptx_text(self, source_path: Path) -> tuple[dict[str, str | int], str]:
        try:
            from pptx import Presentation
        except Exception as exc:
            raise RuntimeError("python-pptx is required for PPTX text extraction.") from exc
        presentation = Presentation(str(source_path))
        lines: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "") or ""
                stripped = text.strip()
                if stripped:
                    slide_lines.append(stripped)
            if slide_lines:
                lines.append(f"Slide {index}")
                lines.extend(slide_lines)
        return {
            "engine": "python-pptx",
            "slide_count": len(presentation.slides),
        }, _normalize_document_text("\n".join(lines))

    def _resolve_office_converter(self) -> str | None:
        configured = self.config.document.office_converter
        if Path(configured).exists():
            return str(Path(configured))
        return shutil.which(configured) or shutil.which("libreoffice")


def _normalize_document_text(text: str) -> str:
    lines = [line.strip() for line in text.splitlines()]
    lines = [line for line in lines if line]
    return "\n".join(lines).strip() + ("\n" if lines else "")
