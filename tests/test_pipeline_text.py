from __future__ import annotations

import json
import shutil
import sys
import uuid
from pathlib import Path
import unittest

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from infoproc.config import load_config
from infoproc.models import DistillMode, ProcessingStage, RunRequest
from infoproc.pipeline import Pipeline


class _StubWriter:
    def __init__(self, title: str) -> None:
        self.title = title

    def write(self, content: str, source_name: str, output_path: Path) -> None:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(f"# {source_name} - {self.title}\n\n{content}", encoding="utf-8")


def _make_pdf(path: Path, text: str) -> None:
    from reportlab.pdfgen import canvas

    path.parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(str(path))
    c.drawString(100, 750, text)
    c.save()


def _make_docx(path: Path, text: str) -> None:
    from docx import Document

    document = Document()
    document.add_paragraph(text)
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(str(path))


def _make_pptx(path: Path, title: str, body: str) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(path))


class PipelineTextTests(unittest.TestCase):
    def test_probe_stage_writes_source_and_probe_only(self) -> None:
        root = Path(__file__).resolve().parents[1] / ".test-temp" / f"probe-{uuid.uuid4().hex}"
        root.mkdir(parents=True, exist_ok=True)
        try:
            input_path = root / "note.txt"
            input_path.write_text("只做探测\n", encoding="utf-8")

            config = load_config(None)
            pipeline = Pipeline(config)
            request = RunRequest(
                input_path=input_path,
                storage_root=root / "storage",
                stage=ProcessingStage.PROBE,
                run_name="probe-run",
            )
            result = pipeline.process(request)

            self.assertEqual(result.files[0].status, "completed")
            run_dir = root / "storage" / "runs" / "probe-run"
            self.assertTrue((run_dir / "00_source" / "plain_text__txt" / "note.txt").exists())
            self.assertTrue((run_dir / "01_probe" / "metadata__json" / "note.json").exists())
            self.assertFalse((run_dir / "04_text_clean" / "clean_text__txt" / "note.txt").exists())
        finally:
            shutil.rmtree(root, ignore_errors=True)

    def test_single_text_input_writes_tree_and_final_outputs(self) -> None:
        root = Path(__file__).resolve().parents[1] / ".test-temp" / f"single-{uuid.uuid4().hex}"
        root.mkdir(parents=True, exist_ok=True)
        try:
            input_path = root / "note.txt"
            input_path.write_text("我们 我们 开始\n我们 我们 开始\n", encoding="utf-8")

            config = load_config(None)
            config.scheduler.llm_workers = 1
            pipeline = Pipeline(config)
            pipeline.distill = _StubWriter("四层知识框架")  # type: ignore[assignment]
            pipeline.rank = _StubWriter("ljg-rank 降秩总结")  # type: ignore[assignment]

            request = RunRequest(
                input_path=input_path,
                storage_root=root / "storage",
                distill_mode=DistillMode.BOTH,
                run_name="single-run",
            )
            result = pipeline.process(request)

            self.assertEqual(result.files[0].status, "completed")
            run_dir = root / "storage" / "runs" / "single-run"
            manifest = json.loads((run_dir / "_manifests" / "run_manifest.json").read_text(encoding="utf-8"))
            self.assertEqual(manifest["status"], "completed")
            self.assertTrue((run_dir / "03_text_raw" / "plain_text__txt" / "note.txt").exists())
            self.assertTrue((run_dir / "04_text_clean" / "clean_text__txt" / "note.txt").exists())
            self.assertTrue((run_dir / "05_final" / "distill__md" / "note.md").exists())
            self.assertTrue((run_dir / "05_final" / "rank__md" / "note.md").exists())
            self.assertTrue((run_dir / "05_final" / "_summaries" / "蒸馏与降秩汇总.md").exists())
            self.assertTrue((run_dir / "05_final" / "_summaries" / "clean汇总.md").exists())
        finally:
            shutil.rmtree(root, ignore_errors=True)

    def test_batch_preserves_relative_structure_and_summaries(self) -> None:
        root = Path(__file__).resolve().parents[1] / ".test-temp" / f"batch-{uuid.uuid4().hex}"
        root.mkdir(parents=True, exist_ok=True)
        try:
            nested = root / "input" / "nested"
            nested.mkdir(parents=True)
            (nested / "a.txt").write_text("第一段\n第一段\n", encoding="utf-8")
            (root / "input" / "b.md").write_text("# 标题\n第二段\n", encoding="utf-8")

            config = load_config(None)
            config.scheduler.llm_workers = 1
            pipeline = Pipeline(config)
            pipeline.distill = _StubWriter("四层知识框架")  # type: ignore[assignment]
            pipeline.rank = _StubWriter("ljg-rank 降秩总结")  # type: ignore[assignment]

            request = RunRequest(
                input_path=root / "input",
                storage_root=root / "storage",
                distill_mode=DistillMode.BOTH,
                recursive=True,
                run_name="batch-run",
            )
            result = pipeline.process(request)

            self.assertEqual(len(result.files), 2)
            run_dir = root / "storage" / "runs" / "batch-run"
            self.assertTrue((run_dir / "00_source" / "plain_text__txt" / "nested" / "a.txt").exists())
            self.assertTrue((run_dir / "00_source" / "markdown__md" / "b.md").exists())
            self.assertTrue((run_dir / "05_final" / "distill__md" / "nested" / "a.md").exists())
            self.assertTrue((run_dir / "05_final" / "rank__md" / "b.md").exists())
            distill_rank_summary = (run_dir / "05_final" / "_summaries" / "蒸馏与降秩汇总.md").read_text(encoding="utf-8")
            clean_summary = (run_dir / "05_final" / "_summaries" / "clean汇总.md").read_text(encoding="utf-8")
            self.assertIn("nested/a.md", distill_rank_summary)
            self.assertIn("b.md", distill_rank_summary)
            self.assertIn("nested/a.txt", clean_summary)
            self.assertIn("b.txt", clean_summary)
        finally:
            shutil.rmtree(root, ignore_errors=True)

    def test_document_extraction_supports_docx_pptx_and_pdf(self) -> None:
        root = Path(__file__).resolve().parents[1] / ".test-temp" / f"docs-{uuid.uuid4().hex}"
        root.mkdir(parents=True, exist_ok=True)
        try:
            input_dir = root / "input"
            _make_docx(input_dir / "guide.docx", "DOCX 内容")
            _make_pptx(input_dir / "slides.pptx", "PPT 标题", "PPT 内容")
            _make_pdf(input_dir / "paper.pdf", "PDF content")

            config = load_config(None)
            pipeline = Pipeline(config)
            request = RunRequest(
                input_path=input_dir,
                storage_root=root / "storage",
                recursive=True,
                stage=ProcessingStage.CLEAN,
                run_name="docs-run",
            )
            result = pipeline.process(request)

            self.assertTrue(all(item.status == "completed" for item in result.files))
            run_dir = root / "storage" / "runs" / "docs-run"
            docx_text = (run_dir / "03_text_raw" / "plain_text__txt" / "guide.txt").read_text(encoding="utf-8")
            pptx_text = (run_dir / "03_text_raw" / "plain_text__txt" / "slides.txt").read_text(encoding="utf-8")
            pdf_text = (run_dir / "03_text_raw" / "plain_text__txt" / "paper.txt").read_text(encoding="utf-8")
            self.assertIn("DOCX 内容", docx_text)
            self.assertIn("PPT 标题", pptx_text)
            self.assertIn("PDF content", pdf_text)
        finally:
            shutil.rmtree(root, ignore_errors=True)

    def test_legacy_office_without_converter_fails_explicitly(self) -> None:
        root = Path(__file__).resolve().parents[1] / ".test-temp" / f"legacy-{uuid.uuid4().hex}"
        root.mkdir(parents=True, exist_ok=True)
        try:
            input_path = root / "legacy.doc"
            input_path.write_bytes(b"fake-doc")

            config = load_config(None)
            config.document.office_converter = "definitely-missing-office-converter"
            pipeline = Pipeline(config)
            request = RunRequest(
                input_path=input_path,
                storage_root=root / "storage",
                stage=ProcessingStage.CLEAN,
                run_name="legacy-run",
            )
            result = pipeline.process(request)

            self.assertEqual(result.files[0].status, "failed")
            self.assertIn("LibreOffice", result.files[0].error or "")
        finally:
            shutil.rmtree(root, ignore_errors=True)


if __name__ == "__main__":
    unittest.main()
